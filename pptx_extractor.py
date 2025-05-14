#!/usr/bin/env python3
import os
import json
import shutil
from pptx import Presentation
from PIL import Image
import io
import uuid

def extract_pptx_content(pptx_path, output_dir):
    """
    Extract content from a PowerPoint file and save it as JSON and images
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        output_dir (str): Directory to save the output files
    
    Returns:
        str: Path to the JSON file
    """
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Create images directory
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)
    
    # Extract filename without extension
    filename = os.path.basename(pptx_path).split('.')[0]
    
    # Load presentation
    prs = Presentation(pptx_path)
    
    # Initialize JSON structure
    course_data = {
        "title": filename,
        "slides": []
    }
    
    # Process each slide
    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slideNumber": i + 1,
            "title": "",
            "content": [],
            "images": []
        }
        
        # Extract text content from slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # If this is likely a title (first text element with content)
                if not slide_data["title"] and shape.text.strip():
                    slide_data["title"] = shape.text.strip()
                else:
                    slide_data["content"].append(shape.text.strip())
            
            # Extract images
            if shape.shape_type == 13:  # 13 is the enum value for pictures
                image = shape.image
                image_bytes = image.blob
                
                # Generate unique filename for the image
                image_filename = f"{filename}_slide_{i+1}_{uuid.uuid4()}.png"
                image_path = os.path.join(images_dir, image_filename)
                
                # Save the image file
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                
                # Add image reference to slide data
                slide_data["images"].append({
                    "filename": image_filename,
                    "path": f"images/{image_filename}"
                })
        
        # Add slide data to course data
        course_data["slides"].append(slide_data)
    
    # Save JSON file
    json_path = os.path.join(output_dir, f"{filename}.json")
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(course_data, f, indent=2, ensure_ascii=False)
    
    return json_path

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Extract content from PowerPoint files')
    parser.add_argument('pptx_path', help='Path to the PowerPoint file')
    parser.add_argument('--output-dir', default='output', help='Directory to save the output files')
    
    args = parser.parse_args()
    
    json_path = extract_pptx_content(args.pptx_path, args.output_dir)
    print(f"Content extracted successfully. JSON saved to: {json_path}") 